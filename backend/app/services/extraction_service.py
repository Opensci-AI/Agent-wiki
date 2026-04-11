import io
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

# MIME types for multimodal extraction
_IMAGE_MIME_TYPES = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "webp": "image/webp",
    "gif": "image/gif",
}


def extract_text(data: bytes, content_type: str) -> str:
    """Synchronous text extraction (no LLM support for images)."""
    try:
        if content_type == "pdf":
            return _extract_pdf(data)
        elif content_type == "docx":
            return _extract_docx(data)
        elif content_type == "pptx":
            return _extract_pptx(data)
        elif content_type in ("xlsx", "xls", "ods"):
            return _extract_xlsx(data)
        elif content_type in ("txt", "md", "csv"):
            return data.decode("utf-8", errors="replace")
        elif content_type in ("png", "jpg", "jpeg", "webp", "gif"):
            return "[Image file — use extract_text_async with LLM config for OCR]"
        else:
            return f"[Unsupported format: {content_type}]"
    except Exception as e:
        return f"[Extraction failed: {str(e)}]"


async def extract_text_async(data: bytes, content_type: str, llm_config: dict | None = None) -> str:
    """Async text extraction with LLM support for images and scanned PDFs.

    Args:
        data: File data bytes
        content_type: File type (pdf, png, jpg, etc.)
        llm_config: LLM configuration for multimodal extraction (required for images)

    Returns:
        Extracted text
    """
    # For non-image types, use sync extraction
    if content_type not in _IMAGE_MIME_TYPES:
        if content_type == "pdf":
            text = _extract_pdf(data)
            # If PDF extraction returned placeholder (scanned PDF), try LLM
            if text.startswith("[PDF appears to be scanned") and llm_config:
                return await _extract_with_llm(data, "application/pdf", llm_config)
            return text
        return extract_text(data, content_type)

    # For images, require LLM config
    if not llm_config:
        return "[Image file — LLM config required for extraction]"

    mime_type = _IMAGE_MIME_TYPES.get(content_type, "image/png")
    return await _extract_with_llm(data, mime_type, llm_config)


async def _extract_with_llm(data: bytes, mime_type: str, llm_config: dict) -> str:
    """Extract text from image/PDF using multimodal LLM."""
    from app.core.llm_client import extract_text_from_image

    prompt = (
        "Extract ALL text visible in this image. "
        "Include headers, paragraphs, tables, labels, and any other text. "
        "Maintain the original structure and formatting as much as possible. "
        "If there are tables, format them with | separators. "
        "Return only the extracted text, no explanations or commentary."
    )

    try:
        return await extract_text_from_image(llm_config, data, mime_type, prompt)
    except Exception as e:
        return f"[LLM extraction failed: {str(e)}]"

def _extract_pdf(data: bytes) -> str:
    text = pdf_extract_text(io.BytesIO(data))
    if text and text.strip():
        return text.strip()
    return "[PDF appears to be scanned/image-based — extraction requires OpenRouter multimodal]"

def _extract_docx(data: bytes) -> str:
    doc = DocxDocument(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)

def _extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides_text)

def _extract_xlsx(data: bytes) -> str:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            sheets_text.append(f"--- {sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets_text)
